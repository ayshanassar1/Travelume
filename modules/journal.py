"""
Journal Module for Travelume
Handles travel journal creation, photo uploads, PPT presentations, and PDF generation
Creates a digital memory book from your travel experiences
"""

import streamlit as st
import os
import uuid
from datetime import datetime
from typing import List, Optional, Dict
from pathlib import Path
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
import json
from PIL import Image as PILImage
import io

# PPT handling
try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Import database
from modules.database import get_database

class JournalManager:
    """Manages travel journals including creation, editing, and PDF generation"""
    
    def __init__(self):
        self.db = get_database()
        self.journals_dir = Path("data/journals")
        self.images_dir = self.journals_dir / "images"
        self.pdfs_dir = self.journals_dir / "pdfs"
        self.ppts_dir = self.journals_dir / "ppts"
        
        # Create directories if they don't exist
        self.journals_dir.mkdir(exist_ok=True, parents=True)
        self.images_dir.mkdir(exist_ok=True, parents=True)
        self.pdfs_dir.mkdir(exist_ok=True, parents=True)
        self.ppts_dir.mkdir(exist_ok=True, parents=True)
    
    def create_journal_ui(self, user_email: str):
        """Display UI for creating a new journal"""
        st.header("📝 Create New Travel Journal")
        
        with st.form("create_journal_form"):
            col1, col2 = st.columns(2)
            
            with col1:
                title = st.text_input("Journal Title*", placeholder="e.g., My Bali Adventure")
                description = st.text_area("Description*", 
                                          placeholder="Share your travel experiences, memories, and stories...",
                                          height=150)
            
            with col2:
                # Get user's trips for linking
                user_trips = self.db.get_user_trips(user_email)
                trip_options = ["None - Standalone Journal"] + [f"{trip.get('destination', 'Unknown')} ({trip.get('id')})" for trip in user_trips]
                selected_trip = st.selectbox("Link to Trip (Optional)", trip_options)
                
                trip_id = None
                if selected_trip != "None - Standalone Journal":
                    # Extract trip ID from selection
                    trip_id = selected_trip.split("(")[-1].replace(")", "")
                
                tags = st.text_input("Tags (comma separated)", 
                                    placeholder="e.g., beach, family, adventure")
                
                is_public = st.checkbox("Make this journal public", value=False)
            
            # Photo upload section
            st.subheader("📸 Add Photos")
            uploaded_photos = st.file_uploader("Upload photos from your trip", 
                                             type=['jpg', 'jpeg', 'png', 'gif'],
                                             accept_multiple_files=True,
                                             help="Upload multiple photos from your travel experience")
            
            # PPT upload section for presentations
            st.subheader("📊 Add Presentation (Optional)")
            if PPTX_AVAILABLE:
                uploaded_ppts = st.file_uploader("Upload PowerPoint presentations", 
                                               type=['pptx', 'ppt'],
                                               accept_multiple_files=True,
                                               help="Upload PPT files with travel photos/memories - slides will be included in your digital memory book")
            else:
                st.info("📌 PPT upload requires python-pptx. Install with: pip install python-pptx")
                uploaded_ppts = []
            
            submitted = st.form_submit_button("✨ Create Digital Memory Book")
            
            if submitted:
                if not title or not description:
                    st.error("Please fill in all required fields (title and description)")
                    return
                
                # Process tags
                tag_list = [tag.strip() for tag in tags.split(",") if tag.strip()] if tags else []
                
                # Process and save uploaded photos
                photo_urls = []
                if uploaded_photos:
                    for uploaded_file in uploaded_photos:
                        try:
                            # Save the uploaded file
                            photo_url = self._save_uploaded_photo(uploaded_file, user_email)
                            if photo_url:
                                photo_urls.append(photo_url)
                        except Exception as e:
                            st.warning(f"Could not save photo {uploaded_file.name}: {str(e)}")
                
                # Process and save uploaded PPT files
                ppt_urls = []
                if uploaded_ppts:
                    for ppt_file in uploaded_ppts:
                        try:
                            ppt_url = self._save_uploaded_ppt(ppt_file, user_email)
                            if ppt_url:
                                ppt_urls.append(ppt_url)
                        except Exception as e:
                            st.warning(f"Could not save PPT {ppt_file.name}: {str(e)}")
                
                # Create journal in database
                journal_id = self.db.create_journal(
                    user_email=user_email,
                    title=title,
                    description=description,
                    photos=photo_urls,
                    trip_id=trip_id
                )
                
                # Update additional fields including PPT files
                updates = {}
                if tag_list:
                    updates['tags'] = tag_list
                if is_public:
                    updates['is_public'] = is_public
                if ppt_urls:
                    updates['ppts'] = ppt_urls
                
                if updates:
                    self.db.update_journal(journal_id, **updates)
                
                st.success("✅ Digital Memory Book created successfully!")
                st.balloons()
                
                # Show next steps
                with st.expander("What's next?"):
                    st.markdown("""
                    1. **Edit Journal**: Add more photos or presentations
                    2. **Generate PDF**: Create a beautiful PDF memory book
                    3. **Share**: Share your travel memories with friends
                    """)
    
    def _save_uploaded_photo(self, uploaded_file, user_email: str) -> str:
        """Save uploaded photo and return its URL/path"""
        try:
            # Create user-specific directory
            user_dir = self.images_dir / user_email.replace("@", "_")
            user_dir.mkdir(exist_ok=True, parents=True)
            
            # Generate unique filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{timestamp}_{uploaded_file.name}"
            filepath = user_dir / filename
            
            # Save the file
            with open(filepath, "wb") as f:
                f.write(uploaded_file.getbuffer())
            
            # Return relative path for storage
            return str(filepath.relative_to("data"))
        except Exception as e:
            st.error(f"Error saving photo: {str(e)}")
            return None
    
    def _save_uploaded_ppt(self, uploaded_file, user_email: str) -> str:
        """Save uploaded PPT file and return its URL/path"""
        try:
            # Create user-specific PPT directory
            user_dir = self.ppts_dir / user_email.replace("@", "_")
            user_dir.mkdir(exist_ok=True, parents=True)
            
            # Generate unique filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{timestamp}_{uploaded_file.name}"
            filepath = user_dir / filename
            
            # Save the file
            with open(filepath, "wb") as f:
                f.write(uploaded_file.getbuffer())
            
            # Return relative path for storage
            return str(filepath.relative_to("data"))
        except Exception as e:
            st.error(f"Error saving PPT: {str(e)}")
            return None
    
    def _extract_slides_from_ppt(self, ppt_path: str, max_slides: int = 20) -> List[str]:
        """Extract images from PPT slides for PDF generation"""
        extracted_images = []
        
        if not PPTX_AVAILABLE:
            return extracted_images
        
        try:
            full_path = os.path.join("data", ppt_path)
            if not os.path.exists(full_path):
                return extracted_images
            
            # Open the presentation
            prs = Presentation(full_path)
            
            # Create temp directory for extracted images
            temp_dir = self.images_dir / "temp_slides"
            temp_dir.mkdir(exist_ok=True, parents=True)
            
            slide_count = 0
            for slide_idx, slide in enumerate(prs.slides):
                if slide_count >= max_slides:
                    break
                
                # Extract images from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            image_bytes = shape.image.blob
                            image_ext = shape.image.ext
                            
                            # Save extracted image
                            img_filename = f"slide_{slide_idx}_{uuid.uuid4().hex[:8]}.{image_ext}"
                            img_path = temp_dir / img_filename
                            
                            with open(img_path, "wb") as f:
                                f.write(image_bytes)
                            
                            extracted_images.append(str(img_path))
                            slide_count += 1
                        except Exception:
                            continue
            
            return extracted_images
        except Exception as e:
            st.warning(f"Could not extract slides from PPT: {str(e)}")
            return extracted_images
    
    def display_user_journals(self, user_email: str):
        """Display all journals for a user"""
        st.header("📚 My Travel Journals")
        
        journals = self.db.get_user_journals(user_email)
        
        if not journals:
            st.info("You haven't created any travel journals yet.")
            st.markdown("""
            Create your first journal to:
            - 📸 Preserve travel photos
            - 📝 Document memories
            - 📄 Generate beautiful PDF keepsakes
            - 🤝 Share experiences with others
            """)
            
            if st.button("Create Your First Journal", type="primary"):
                st.session_state.journal_action = "create"
                st.rerun()
            return
        
        # Search and filter options
        col1, col2, col3 = st.columns([2, 2, 1])
        
        with col1:
            search_query = st.text_input("🔍 Search journals", placeholder="Search by title or tags")
        
        with col2:
            sort_by = st.selectbox("Sort by", ["Newest", "Oldest", "Title A-Z", "Title Z-A"])
        
        with col3:
            show_only_pdf = st.checkbox("PDF Ready", help="Show only journals with PDF generated")
        
        # Filter journals
        filtered_journals = journals
        
        if search_query:
            filtered_journals = [
                j for j in filtered_journals 
                if search_query.lower() in j.get("title", "").lower() 
                or search_query.lower() in " ".join(j.get("tags", [])).lower()
            ]
        
        if show_only_pdf:
            filtered_journals = [j for j in filtered_journals if j.get("pdf_generated", False)]
        
        # Sort journals
        if sort_by == "Newest":
            filtered_journals.sort(key=lambda x: x.get("created_at", ""), reverse=True)
        elif sort_by == "Oldest":
            filtered_journals.sort(key=lambda x: x.get("created_at", ""))
        elif sort_by == "Title A-Z":
            filtered_journals.sort(key=lambda x: x.get("title", "").lower())
        elif sort_by == "Title Z-A":
            filtered_journals.sort(key=lambda x: x.get("title", "").lower(), reverse=True)
        
        # Display journals
        for journal in filtered_journals:
            self._display_journal_card(journal)
    
    def _display_journal_card(self, journal: Dict):
        """Display a journal card in the list view"""
        journal_id = journal.get("id")
        title = journal.get("title", "Untitled Journal")
        description = journal.get("description", "")
        created_date = journal.get("created_at", "")
        photos = journal.get("photos", [])
        tags = journal.get("tags", [])
        has_pdf = journal.get("pdf_generated", False)
        
        with st.container():
            col1, col2 = st.columns([1, 3])
            
            with col1:
                # Display first photo as thumbnail or placeholder
                if photos:
                    try:
                        photo_path = os.path.join("data", photos[0])
                        if os.path.exists(photo_path):
                            st.image(photo_path, use_column_width=True)
                    except:
                        st.image("https://images.unsplash.com/photo-1488646953014-85cb44e25828?w=400", 
                                caption="Travel Memories")
                else:
                    st.image("https://images.unsplash.com/photo-1488646953014-85cb44e25828?w-400", 
                            caption="No photos yet")
            
            with col2:
                st.subheader(title)
                st.caption(f"Created: {created_date}")
                
                # Display description preview
                if description:
                    preview = description[:200] + "..." if len(description) > 200 else description
                    st.write(preview)
                
                # Display tags
                if tags:
                    tag_html = " ".join([f"<span style='background-color: #e0f2fe; padding: 4px 8px; border-radius: 12px; margin-right: 5px; font-size: 0.8em;'>{tag}</span>" for tag in tags])
                    st.markdown(tag_html, unsafe_allow_html=True)
                
                # Action buttons
                col_btn1, col_btn2, col_btn3, col_btn4 = st.columns(4)
                
                with col_btn1:
                    if st.button("📖 View", key=f"view_{journal_id}", use_container_width=True):
                        st.session_state.current_journal = journal_id
                        st.session_state.journal_action = "view"
                        st.rerun()
                
                with col_btn2:
                    if st.button("✏️ Edit", key=f"edit_{journal_id}", use_container_width=True):
                        st.session_state.current_journal = journal_id
                        st.session_state.journal_action = "edit"
                        st.rerun()
                
                with col_btn3:
                    if has_pdf:
                        pdf_path = journal.get("pdf_path")
                        if pdf_path and os.path.exists(pdf_path):
                            with open(pdf_path, "rb") as pdf_file:
                                st.download_button(
                                    label="📄 PDF",
                                    data=pdf_file,
                                    file_name=f"{title.replace(' ', '_')}.pdf",
                                    mime="application/pdf",
                                    key=f"pdf_{journal_id}",
                                    use_container_width=True
                                )
                        else:
                            st.button("📄 Generate", key=f"gen_{journal_id}", 
                                     disabled=True, use_container_width=True)
                    else:
                        if st.button("📄 Generate", key=f"gen_{journal_id}", use_container_width=True):
                            st.session_state.current_journal = journal_id
                            st.session_state.journal_action = "generate_pdf"
                            st.rerun()
                
                with col_btn4:
                    if st.button("🗑️ Delete", key=f"delete_{journal_id}", 
                                type="secondary", use_container_width=True):
                        if st.session_state.get("delete_confirm") == journal_id:
                            if self.db.delete_journal(st.session_state.user_email, journal_id):
                                st.success("Journal deleted!")
                                st.rerun()
                        else:
                            st.session_state.delete_confirm = journal_id
                            st.warning("Click delete again to confirm")
                            st.rerun()
            
            st.divider()
    
    def view_journal(self, journal_id: str, user_email: str):
        """Display a full journal view"""
        journal = self.db.get_journal(journal_id)
        
        if not journal:
            st.error("Journal not found!")
            if st.button("Back to Journals"):
                st.session_state.journal_action = None
                st.rerun()
            return
        
        # Update view count
        self.db.update_journal(journal_id, views=journal.get("views", 0) + 1)
        
        # Back button
        if st.button("← Back to Journals"):
            st.session_state.journal_action = None
            st.rerun()
        
        # Journal header
        col_title, col_actions = st.columns([3, 1])
        
        with col_title:
            st.title(journal.get("title", "Untitled Journal"))
        
        with col_actions:
            if st.button("✏️ Edit Journal"):
                st.session_state.journal_action = "edit"
                st.rerun()
        
        # Metadata
        meta_col1, meta_col2, meta_col3, meta_col4 = st.columns(4)
        with meta_col1:
            st.metric("Created", journal.get("created_at", "").split()[0])
        with meta_col2:
            st.metric("Updated", journal.get("updated_at", "").split()[0])
        with meta_col3:
            st.metric("Views", journal.get("views", 0))
        with meta_col4:
            st.metric("Photos", len(journal.get("photos", [])))
        
        # Linked trip info
        if journal.get("trip_id"):
            trip = self.db.get_trip_details(journal.get("trip_id"))
            if trip:
                st.info(f"📌 Linked to trip: **{trip.get('destination', 'Unknown')}**")
        
        # Tags
        tags = journal.get("tags", [])
        if tags:
            tag_html = " ".join([f"<span style='background-color: #3b82f6; color: white; padding: 4px 12px; border-radius: 20px; margin-right: 8px; font-size: 0.9em;'>{tag}</span>" for tag in tags])
            st.markdown(tag_html, unsafe_allow_html=True)
        
        st.divider()
        
        # Description
        st.subheader("My Travel Story")
        st.write(journal.get("description", ""))
        
        # Photos section
        photos = journal.get("photos", [])
        if photos:
            st.subheader(f"📸 Travel Photos ({len(photos)})")
            
            # Display photos in a grid
            cols = st.columns(min(3, len(photos)))
            for idx, photo_path in enumerate(photos):
                full_path = os.path.join("data", photo_path)
                if os.path.exists(full_path):
                    with cols[idx % 3]:
                        st.image(full_path, use_column_width=True)
                        st.caption(f"Photo {idx + 1}")
                else:
                    st.warning(f"Photo not found: {photo_path}")
        else:
            st.info("No photos added yet. Add photos to make your journal more memorable!")
        
        # PPT Presentations section
        ppts = journal.get("ppts", [])
        if ppts:
            st.subheader(f"📊 Presentations ({len(ppts)})")
            
            for idx, ppt_path in enumerate(ppts):
                full_path = os.path.join("data", ppt_path)
                ppt_name = os.path.basename(ppt_path)
                
                col_ppt1, col_ppt2 = st.columns([3, 1])
                with col_ppt1:
                    st.markdown(f"📎 **{ppt_name}**")
                with col_ppt2:
                    if os.path.exists(full_path):
                        with open(full_path, "rb") as ppt_file:
                            st.download_button(
                                label="Download",
                                data=ppt_file,
                                file_name=ppt_name,
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                key=f"download_ppt_{idx}",
                                use_container_width=True
                            )
                    else:
                        st.warning("File not found")
        
        # PDF section
        st.divider()
        if journal.get("pdf_generated"):
            pdf_path = journal.get("pdf_path")
            if pdf_path and os.path.exists(pdf_path):
                st.subheader("📄 PDF Version")
                col_pdf1, col_pdf2 = st.columns([3, 1])
                with col_pdf1:
                    st.success("PDF has been generated! Download your travel keepsake.")
                with col_pdf2:
                    with open(pdf_path, "rb") as pdf_file:
                        st.download_button(
                            label="Download PDF",
                            data=pdf_file,
                            file_name=f"{journal.get('title', 'journal').replace(' ', '_')}.pdf",
                            mime="application/pdf",
                            use_container_width=True
                        )
            else:
                if st.button("🔄 Regenerate PDF", type="secondary"):
                    self._generate_pdf(journal_id, user_email)
        else:
            st.subheader("Create PDF Keepsake")
            st.info("Generate a beautiful PDF version of your journal to save or share!")
            if st.button("✨ Generate PDF Now", type="primary"):
                with st.spinner("Creating your PDF keepsake..."):
                    success = self._generate_pdf(journal_id, user_email)
                    if success:
                        st.success("PDF generated successfully!")
                        st.rerun()
                    else:
                        st.error("Failed to generate PDF. Please try again.")
    
    def edit_journal(self, journal_id: str, user_email: str):
        """Edit an existing journal"""
        journal = self.db.get_journal(journal_id)
        
        if not journal:
            st.error("Journal not found!")
            return
        
        st.header("✏️ Edit Journal")
        
        with st.form("edit_journal_form"):
            title = st.text_input("Title", value=journal.get("title", ""))
            description = st.text_area("Description", 
                                      value=journal.get("description", ""),
                                      height=200)
            
            # Tags
            current_tags = ", ".join(journal.get("tags", []))
            tags = st.text_input("Tags (comma separated)", 
                                value=current_tags,
                                placeholder="e.g., beach, family, adventure")
            
            is_public = st.checkbox("Make this journal public", 
                                   value=journal.get("is_public", False))
            
            # Photo management
            st.subheader("📸 Manage Photos")
            current_photos = journal.get("photos", [])
            
            # Display current photos
            if current_photos:
                st.write("Current photos:")
                cols = st.columns(min(4, len(current_photos)))
                for idx, photo_path in enumerate(current_photos):
                    full_path = os.path.join("data", photo_path)
                    if os.path.exists(full_path):
                        with cols[idx % 4]:
                            st.image(full_path, use_column_width=True)
                            if st.button(f"Remove {idx + 1}", key=f"remove_{idx}"):
                                # Remove photo
                                updated_photos = [p for p in current_photos if p != photo_path]
                                self.db.update_journal(journal_id, photos=updated_photos)
                                st.success("Photo removed!")
                                st.rerun()
            
            # Upload new photos
            new_photos = st.file_uploader("Add more photos", 
                                         type=['jpg', 'jpeg', 'png', 'gif'],
                                         accept_multiple_files=True,
                                         help="Upload additional photos")
            
            # PPT management section
            st.subheader("📊 Manage Presentations")
            current_ppts = journal.get("ppts", [])
            
            # Display current PPTs
            if current_ppts:
                st.write("Current presentations:")
                for idx, ppt_path in enumerate(current_ppts):
                    ppt_name = os.path.basename(ppt_path)
                    col_ppt1, col_ppt2 = st.columns([3, 1])
                    with col_ppt1:
                        st.markdown(f"📎 {ppt_name}")
                    with col_ppt2:
                        if st.button(f"Remove", key=f"remove_ppt_{idx}"):
                            updated_ppts = [p for p in current_ppts if p != ppt_path]
                            self.db.update_journal(journal_id, ppts=updated_ppts)
                            st.success("Presentation removed!")
                            st.rerun()
            
            # Upload new PPTs
            if PPTX_AVAILABLE:
                new_ppts = st.file_uploader("Add presentations", 
                                           type=['pptx', 'ppt'],
                                           accept_multiple_files=True,
                                           help="Upload additional PowerPoint files")
            else:
                st.info("📌 PPT upload requires python-pptx")
                new_ppts = []
            
            col1, col2 = st.columns(2)
            with col1:
                submit = st.form_submit_button("💾 Save Changes", type="primary")
            with col2:
                cancel = st.form_submit_button("❌ Cancel")
            
            if cancel:
                st.session_state.journal_action = "view"
                st.rerun()
            
            if submit:
                # Process tags
                tag_list = [tag.strip() for tag in tags.split(",") if tag.strip()] if tags else []
                
                # Process new photos
                new_photo_urls = []
                if new_photos:
                    for uploaded_file in new_photos:
                        try:
                            photo_url = self._save_uploaded_photo(uploaded_file, user_email)
                            if photo_url:
                                new_photo_urls.append(photo_url)
                        except Exception as e:
                            st.warning(f"Could not save photo {uploaded_file.name}: {str(e)}")
                
                # Process new PPTs
                new_ppt_urls = []
                if new_ppts:
                    for ppt_file in new_ppts:
                        try:
                            ppt_url = self._save_uploaded_ppt(ppt_file, user_email)
                            if ppt_url:
                                new_ppt_urls.append(ppt_url)
                        except Exception as e:
                            st.warning(f"Could not save PPT {ppt_file.name}: {str(e)}")
                
                # Combine old and new files
                all_photos = current_photos + new_photo_urls
                all_ppts = current_ppts + new_ppt_urls
                
                # Update journal
                updates = {
                    "title": title,
                    "description": description,
                    "tags": tag_list,
                    "photos": all_photos,
                    "ppts": all_ppts,
                    "is_public": is_public
                }
                
                if self.db.update_journal(journal_id, **updates):
                    st.success("✅ Journal updated successfully!")
                    
                    # Mark PDF as outdated if content changed
                    if journal.get("pdf_generated"):
                        self.db.update_journal(journal_id, pdf_generated=False)
                    
                    st.balloons()
                    
                    # Option to go back to view
                    if st.button("View Updated Journal"):
                        st.session_state.journal_action = "view"
                        st.rerun()
    
    def _generate_pdf(self, journal_id: str, user_email: str) -> bool:
        """Generate a beautiful PDF memory book from a journal"""
        try:
            journal = self.db.get_journal(journal_id)
            if not journal:
                return False
            
            # Create PDF filename
            title_safe = "".join(c if c.isalnum() else "_" for c in journal.get("title", "journal"))
            pdf_filename = f"{title_safe}_{journal_id}.pdf"
            pdf_path = self.pdfs_dir / pdf_filename
            
            # Setup Document
            doc = SimpleDocTemplate(
                str(pdf_path), 
                pagesize=letter,
                rightMargin=50, leftMargin=50, 
                topMargin=50, bottomMargin=50
            )
            
            story = []
            styles = getSampleStyleSheet()
            
            # Custom Styles
            styles.add(ParagraphStyle(name='JournalTitle', parent=styles['Heading1'], fontSize=24, spaceAfter=20, textColor=colors.HexColor("#831843")))
            styles.add(ParagraphStyle(name='JournalText', parent=styles['Normal'], fontSize=12, leading=16, spaceBefore=10, textColor=colors.HexColor("#374151")))
            styles.add(ParagraphStyle(name='Caption', parent=styles['Italic'], fontSize=10, textColor=colors.gray, alignment=1)) # Center align
            
            # --- Cover Page Logic ---
            def create_cover_page(canvas, doc):
                """Draws a cover page with a background color/pattern and centered title"""
                canvas.saveState()
                
                # Background
                canvas.setFillColor(colors.HexColor("#fce7f3")) # Light pink/rose background
                canvas.rect(0, 0, letter[0], letter[1], fill=1, stroke=0)
                
                # Decorative border
                canvas.setStrokeColor(colors.HexColor("#be185d"))
                canvas.setLineWidth(3)
                canvas.rect(0.5*inch, 0.5*inch, letter[0]-1*inch, letter[1]-1*inch)
                
                # Title
                canvas.setFont('Helvetica-Bold', 36)
                canvas.setFillColor(colors.HexColor("#831843"))
                # Center text logic
                text_width = canvas.stringWidth(journal.get("title", "Travel Journal"), 'Helvetica-Bold', 36)
                canvas.drawString((letter[0] - text_width) / 2.0, letter[1]/2.0 + 50, journal.get("title", "Travel Journal"))
                
                # Subtitle
                username = user_email.split("@")[0]
                date_str = datetime.now().strftime("%B %Y")
                subtitle = f"Created by {username} • {date_str}"
                
                canvas.setFont('Helvetica', 18)
                canvas.setFillColor(colors.HexColor("#9d174d"))
                text_width = canvas.stringWidth(subtitle, 'Helvetica', 18)
                canvas.drawString((letter[0] - text_width) / 2.0, letter[1]/2.0 - 20, subtitle)
                
                # Add a quote or decorative text
                canvas.setFont('Helvetica-Oblique', 14)
                canvas.setFillColor(colors.HexColor("#be185d"))
                quote = "Every journey tells a story..."
                text_width = canvas.stringWidth(quote, 'Helvetica-Oblique', 14)
                canvas.drawString((letter[0] - text_width) / 2.0, letter[1]/2.0 - 80, quote)

                canvas.restoreState()

            # --- Content Pages ---
            
            # Title Page Header
            story.append(Paragraph("My Travel Story", styles['JournalTitle']))
            story.append(Spacer(1, 10))
            
            if journal.get("description"):
                story.append(Paragraph(journal.get("description", ""), styles['JournalText']))
                story.append(Spacer(1, 20))
            
            # Photos Grid
            photos = journal.get("photos", [])
            if photos:
                story.append(Paragraph("Memories", styles['Heading2']))
                story.append(Spacer(1, 10))
                
                # Prepare rows of 2 photos
                photo_rows = []
                current_row = []
                
                for photo_path in photos:
                    # Handle both relative paths from 'data/' and backward compatibility
                    # The db usually stores "journals/images/..." which is relative to "data/"
                    # But reportlab needs absolute path or relative to CWD
                    
                    # Try to find the file
                    possible_paths = [
                        Path("data") / photo_path,
                        Path(photo_path)
                    ]
                    
                    full_path = None
                    for p in possible_paths:
                        if p.exists():
                            full_path = p
                            break
                    
                    if full_path:
                        try:
                            # Resize/Aspect Ratio check could go here
                            # Maintain aspect ratio logic could be added
                            img = Image(str(full_path), width=3*inch, height=2.25*inch)
                            current_row.append(img)
                        except Exception:
                            pass
                    
                    if len(current_row) == 2:
                        photo_rows.append(current_row)
                        current_row = []
                
                if current_row:
                    photo_rows.append(current_row)
                    
                # Create Table
                if photo_rows:
                    # Calculate row heights dynamically or set fixed
                    t = Table(photo_rows, colWidths=[3.2*inch, 3.2*inch])
                    t.setStyle(TableStyle([
                        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
                        ('VALIGN', (0,0), (-1,-1), 'TOP'),
                        ('BOTTOMPADDING', (0,0), (-1,-1), 15),
                        ('LEFTPADDING', (0,0), (-1,-1), 5),
                        ('RIGHTPADDING', (0,0), (-1,-1), 5),
                    ]))
                    story.append(t)
            
            # Helper to add cover page logic
            # reportlab's build method executes onFirstPage for the first page
            # To ensure the cover is its own page, we can use a custom PageTemplate or a simpler trick
            # Simpler trick: Use a custom onFirstPage function that draws the cover, then we PageBreak immediately?
            # Actually, standard doc.build allows onFirstPage=func.
            
            doc.build(story, onFirstPage=create_cover_page)
            
            # Update journal with PDF info (store relative to data/ for consistency)
            # The pdf_path object is data/journals/pdfs/filename.pdf
            # The DB expects path relative to "data" folder usually? 
            # Let's check existing DB logic. 
            # existing _save_uploaded_photo returns "journals/images/..." (relative to data root? No, relative to data)
            # line 178: return str(filepath.relative_to("data"))
            # So we should store relative to data.
            
            rel_path = str(pdf_path.relative_to("data"))
            self.db.mark_pdf_generated(journal_id, rel_path)
            
            return True
            
        except Exception as e:
            # import traceback
            # traceback.print_exc()
            st.error(f"Error generating PDF: {str(e)}")
            return False

def show_journals_page(user_email: str):
    """Main function to show the journals page"""
    manager = JournalManager()
    
    # Initialize session state for journal actions
    if "journal_action" not in st.session_state:
        st.session_state.journal_action = None
    
    if "current_journal" not in st.session_state:
        st.session_state.current_journal = None
    
    # Show appropriate view based on action
    action = st.session_state.journal_action
    
    if action == "create":
        manager.create_journal_ui(user_email)
    
    elif action == "view" and st.session_state.current_journal:
        manager.view_journal(st.session_state.current_journal, user_email)
    
    elif action == "edit" and st.session_state.current_journal:
        manager.edit_journal(st.session_state.current_journal, user_email)
    
    elif action == "generate_pdf" and st.session_state.current_journal:
        with st.spinner("Generating PDF..."):
            success = manager._generate_pdf(st.session_state.current_journal, user_email)
            if success:
                st.success("PDF generated successfully!")
                st.session_state.journal_action = "view"
                st.rerun()
            else:
                st.error("Failed to generate PDF")
                st.session_state.journal_action = "view"
    
    else:
        # Default: Show all journals
        manager.display_user_journals(user_email)

# Test function
def test_journal_module():
    """Test the journal module"""
    st.title("Journal Module Test")
    
    # Create a test user if needed
    db = get_database()
    
    if "test_user" not in st.session_state:
        st.session_state.test_user = "test@example.com"
    
    user_email = st.session_state.test_user
    
    # Show journals page
    show_journals_page(user_email)

if __name__ == "__main__":
    test_journal_module()